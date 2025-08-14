from pptx import Presentation
from pptx.util import Inches, Cm
import re
from pptx.dml.color import RGBColor

def sort_files_numerically(file_list):
    """
    Trie une liste de fichiers contenant des nombres dans leurs noms.
    """
    return sorted(file_list, key=lambda x: int(re.search(r'\d+', x).group()))

def create_presentation(image_folder, output_pptx):
    """
    Crée une présentation PowerPoint avec des images en haut à gauche, occupant toute la hauteur.
    
    :param image_folder: Chemin vers le dossier contenant les images.
    :param output_pptx: Nom du fichier PowerPoint de sortie.
    """
    import os

    # Crée une présentation PowerPoint
    prs = Presentation()

    # Dimensions de la diapositive (en pouces)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Récupère la liste des images dans le dossier
 # Récupère la liste des images dans le dossier
    images = [f for f in os.listdir(image_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]

    images2 = [f for f in os.listdir("imageextracted") if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    images2 = sort_files_numerically(images2)  # Trie les fichiers numériquement    # images.sort()  # Tri pour garantir l'ordre



    images = sort_files_numerically(images)  # Trie les fichiers numériquement    # images.sort()  # Tri pour garantir l'ordre

    for i in range(len(images)):
        image_file = images[i]

        if i%100==0 and i<=5000:
            images2files = images2[i//100]
            im2path = os.path.join("imageextracted", images2files)

            # Ajoute une diapositive vierge
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout vide
            # Définir un arrière-plan avec la couleur #e8e8e8
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(232, 232, 232)  # Couleur #e8e8e8
            # Chemin complet de l'image
            image_path = os.path.join(image_folder, image_file)

            # Ajoute l'image à la diapositive
            left = Cm(0)  # Position horizontale (0 cm depuis le bord gauche)
            top = Cm(0)   # Position verticale (0 cm depuis le bord supérieur)
            height = slide_height  # Hauteur de l'image = hauteur de la diapositive

            # Ajoute l'image redimensionnée
            slide.shapes.add_picture(image_path, left, top, height=height)

            # Calculer la position horizontale (à partir du bord droit)
            top = Cm(0)  # Position verticale (0 cm depuis le bord supérieur)
            height = slide_height  # Hauteur de l'image = hauteur de la diapositive
            width = 1750 * height/1000
            #CA MARCHE PAS
            left = Cm(slide_width) - Cm(width/37.795)  # Position horizontale : bord droit - largeur de l'image
            print(left)
            print(slide_width)
            # Ajouter l'image redimensionnée dans le coin en haut à droite
            slide.shapes.add_picture(im2path, left, top, height=height)


    # Sauvegarde la présentation
    prs.save(output_pptx)
    print(f"Présentation créée : {output_pptx}")
# Exemple d'utilisation
image_folder = "tirreussi"  # Remplacez par le chemin de votre dossier
output_pptx = "stop_motion.pptx"
create_presentation(image_folder, output_pptx)